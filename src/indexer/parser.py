"""Document parser for extracting text from various file formats."""

import io
import logging
from pathlib import Path
from typing import Protocol

logger = logging.getLogger(__name__)


class ParserProtocol(Protocol):
    """Protocol for document parsers."""

    def parse(self, content: bytes, file_name: str) -> str:
        """Parse document content to plain text."""
        ...


class TextParser:
    """Parser for plain text files."""

    ENCODINGS = ["utf-8", "utf-16", "latin-1", "cp1252"]

    def parse(self, content: bytes, file_name: str) -> str:
        """Parse text file content."""
        for encoding in self.ENCODINGS:
            try:
                return content.decode(encoding)
            except UnicodeDecodeError:
                continue

        # Fallback: decode with errors replaced
        return content.decode("utf-8", errors="replace")


class MarkdownParser(TextParser):
    """Parser for Markdown files - inherits from TextParser."""

    pass


class PDFParser:
    """Parser for PDF files."""

    def parse(self, content: bytes, file_name: str) -> str:
        """Parse PDF content to plain text."""
        try:
            from pypdf import PdfReader

            reader = PdfReader(io.BytesIO(content))
            text_parts = []

            for page in reader.pages:
                try:
                    text = page.extract_text()
                    if text:
                        text_parts.append(text)
                except Exception as e:
                    logger.warning(f"Error extracting text from PDF page: {e}")
                    continue

            return "\n\n".join(text_parts)

        except ImportError:
            logger.error("pypdf not installed. Cannot parse PDF files.")
            return ""
        except Exception as e:
            logger.error(f"Error parsing PDF {file_name}: {e}")
            return ""


class DocxParser:
    """Parser for DOCX files."""

    def parse(self, content: bytes, file_name: str) -> str:
        """Parse DOCX content to plain text."""
        try:
            from docx import Document

            doc = Document(io.BytesIO(content))
            text_parts = []

            for para in doc.paragraphs:
                if para.text.strip():
                    text_parts.append(para.text)

            # Also extract text from tables
            for table in doc.tables:
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        if cell.text.strip():
                            row_text.append(cell.text.strip())
                    if row_text:
                        text_parts.append(" | ".join(row_text))

            return "\n\n".join(text_parts)

        except ImportError:
            logger.error("python-docx not installed. Cannot parse DOCX files.")
            return ""
        except Exception as e:
            logger.error(f"Error parsing DOCX {file_name}: {e}")
            return ""


class XlsxParser:
    """Parser for XLSX files."""

    def parse(self, content: bytes, file_name: str) -> str:
        """Parse XLSX content to plain text."""
        try:
            from openpyxl import load_workbook

            wb = load_workbook(io.BytesIO(content), read_only=True, data_only=True)
            text_parts = []

            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                text_parts.append(f"=== Sheet: {sheet_name} ===")

                for row in sheet.iter_rows():
                    row_values = []
                    for cell in row:
                        if cell.value is not None:
                            row_values.append(str(cell.value))
                    if row_values:
                        text_parts.append(" | ".join(row_values))

            wb.close()
            return "\n".join(text_parts)

        except ImportError:
            logger.error("openpyxl not installed. Cannot parse XLSX files.")
            return ""
        except Exception as e:
            logger.error(f"Error parsing XLSX {file_name}: {e}")
            return ""


class PptxParser:
    """Parser for PPTX files."""

    def parse(self, content: bytes, file_name: str) -> str:
        """Parse PPTX content to plain text."""
        try:
            from pptx import Presentation

            prs = Presentation(io.BytesIO(content))
            text_parts = []

            for slide_num, slide in enumerate(prs.slides, 1):
                slide_texts = [f"=== Slide {slide_num} ==="]

                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_texts.append(shape.text)

                if len(slide_texts) > 1:  # Has content besides header
                    text_parts.extend(slide_texts)

            return "\n\n".join(text_parts)

        except ImportError:
            logger.error("python-pptx not installed. Cannot parse PPTX files.")
            return ""
        except Exception as e:
            logger.error(f"Error parsing PPTX {file_name}: {e}")
            return ""


class DocumentParser:
    """Main document parser that delegates to format-specific parsers."""

    # Map file extensions to parsers
    PARSERS: dict[str, ParserProtocol] = {
        # Text formats
        ".txt": TextParser(),
        ".md": MarkdownParser(),
        ".markdown": MarkdownParser(),
        ".rst": TextParser(),
        ".csv": TextParser(),
        ".json": TextParser(),
        ".xml": TextParser(),
        ".html": TextParser(),
        ".htm": TextParser(),
        # Office formats
        ".pdf": PDFParser(),
        ".docx": DocxParser(),
        ".doc": DocxParser(),  # May not work for old .doc format
        ".xlsx": XlsxParser(),
        ".xls": XlsxParser(),  # May not work for old .xls format
        ".pptx": PptxParser(),
    }

    SUPPORTED_EXTENSIONS = list(PARSERS.keys())

    @classmethod
    def is_supported(cls, file_path: str | Path) -> bool:
        """Check if a file type is supported."""
        ext = Path(file_path).suffix.lower()
        return ext in cls.PARSERS

    @classmethod
    def parse(cls, content: bytes, file_path: str | Path) -> str:
        """Parse document content to plain text.

        Args:
            content: Raw file content as bytes.
            file_path: Path to the file (used for extension detection).

        Returns:
            Extracted plain text content.
        """
        path = Path(file_path)
        ext = path.suffix.lower()

        parser = cls.PARSERS.get(ext)
        if parser is None:
            logger.warning(f"No parser available for extension: {ext}")
            return ""

        try:
            text = parser.parse(content, path.name)
            # Clean up the text
            return cls._clean_text(text)
        except Exception as e:
            logger.error(f"Error parsing {file_path}: {e}")
            return ""

    @staticmethod
    def _clean_text(text: str) -> str:
        """Clean and normalize extracted text."""
        if not text:
            return ""

        # Replace multiple newlines with double newline
        import re
        text = re.sub(r"\n{3,}", "\n\n", text)

        # Replace multiple spaces with single space
        text = re.sub(r" {2,}", " ", text)

        # Strip leading/trailing whitespace
        text = text.strip()

        return text

    @classmethod
    def get_supported_types(cls) -> list[str]:
        """Get list of supported file extensions."""
        return cls.SUPPORTED_EXTENSIONS.copy()
